"""
Document parser module with Docling integration and fallback loaders.
Supports PDF, Word, Excel, PPT, PNG, and TXT formats.
"""
import sys
from pathlib import Path
from typing import Dict, Any, Tuple

sys.path.insert(0, str(Path(__file__).parent.parent))
from utils import get_logger

logger = get_logger(__name__)


class DocumentParser:
    """
    Document parser using Docling for advanced parsing with fallback to basic loaders.
    Supports PDF, Word, Excel, PPT, PNG, and TXT formats.
    """
    
    def __init__(self):
        self.docling_available = False
        self.converter = None
        
        try:
            from docling.document_converter import DocumentConverter
            self.converter = DocumentConverter()
            self.docling_available = True
            logger.info("Docling initialized successfully")
        except ImportError as e:
            logger.warning(f"Docling not available, using fallback loaders: {e}")
        except Exception as e:
            logger.warning(f"Failed to initialize Docling: {e}")
    
    def parse_document(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """
        Parse a document and extract text content with metadata.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Tuple of (text_content, metadata)
        """
        file_path = Path(file_path)
        suffix = file_path.suffix.lower()
        metadata = {
            "source": str(file_path),
            "filename": file_path.name,
            "file_type": suffix,
            "tables": [],
            "images": []
        }
        
        if self.docling_available and suffix in ['.pdf', '.docx', '.pptx', '.xlsx', '.png', '.jpg', '.jpeg', '.tiff']:
            try:
                return self._parse_with_docling(file_path, metadata)
            except Exception as e:
                logger.warning(f"Docling parsing failed for {file_path}, using fallback: {e}")
        
        return self._parse_with_fallback(file_path, metadata)
    
    def _parse_with_docling(self, file_path: Path, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """Parse document using Docling."""
        logger.info(f"Parsing {file_path} with Docling...")
        
        result = self.converter.convert(str(file_path))
        doc = result.document
        
        markdown_content = doc.export_to_markdown()
        
        tables = []
        images = []
        table_texts = []
        
        try:
            for item, level in doc.iterate_items():
                label = getattr(item, "label", None)
                if label == "table":
                    if hasattr(item, "export_to_dataframe"):
                        df = item.export_to_dataframe()
                        csv_content = df.to_csv(index=False)
                        tables.append({
                            "data": df.to_dict(),
                            "csv": csv_content
                        })
                        
                        table_text = self._format_table_for_indexing(df)
                        if table_text:
                            table_texts.append(table_text)
                elif label in ["figure", "picture", "image"]:
                    images.append({
                        "alt_text": getattr(item, "caption", ""),
                        "description": str(item)
                    })
        except Exception as e:
            logger.warning(f"Error extracting tables/images: {e}")
        
        if table_texts:
            markdown_content += "\n\n--- TABLE DATA ---\n\n" + "\n\n".join(table_texts)
        
        metadata["tables"] = tables
        metadata["images"] = images
        metadata["parser"] = "docling"
        metadata["table_count"] = len(tables)
        
        logger.info(f"✓ Docling parsed {file_path}: {len(markdown_content)} chars, {len(tables)} tables, {len(images)} images")
        return markdown_content, metadata
    
    def _format_table_for_indexing(self, df) -> str:
        """
        Format a DataFrame table into searchable text that preserves row relationships.
        Each row is formatted as a complete record for better retrieval.
        """
        try:
            if df.empty:
                return ""
            
            columns = list(df.columns)
            rows_text = []
            
            for idx, row in df.iterrows():
                row_parts = []
                for col in columns:
                    value = row[col]
                    if value is not None and str(value).strip():
                        row_parts.append(f"{col}: {value}")
                if row_parts:
                    rows_text.append(" | ".join(row_parts))
            
            if rows_text:
                header = f"Table columns: {', '.join(columns)}"
                return header + "\n" + "\n".join(rows_text)
            return ""
        except Exception as e:
            logger.warning(f"Error formatting table: {e}")
            return ""
    
    def _parse_with_fallback(self, file_path: Path, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """Parse document using fallback loaders."""
        suffix = file_path.suffix.lower()
        content = ""
        
        try:
            if suffix == '.pdf':
                content = self._parse_pdf_fallback(file_path)
            elif suffix in ['.docx', '.doc']:
                content = self._parse_word_fallback(file_path)
            elif suffix in ['.xlsx', '.xls']:
                content = self._parse_excel_fallback(file_path)
            elif suffix in ['.pptx', '.ppt']:
                content = self._parse_ppt_fallback(file_path)
            elif suffix == '.txt':
                content = self._parse_txt_fallback(file_path)
            elif suffix in ['.png', '.jpg', '.jpeg', '.tiff']:
                content = self._parse_image_fallback(file_path)
            else:
                content = f"Unsupported file format: {suffix}"
                logger.warning(content)
        except Exception as e:
            content = f"Error parsing file: {str(e)}"
            logger.error(content)
        
        metadata["parser"] = "fallback"
        logger.info(f"✓ Fallback parsed {file_path}: {len(content)} chars")
        return content, metadata
    
    def _parse_pdf_fallback(self, file_path: Path) -> str:
        """Parse PDF using pypdf with improved table detection."""
        from pypdf import PdfReader
        reader = PdfReader(str(file_path))
        text_parts = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            processed_text = self._process_pdf_text_for_tables(text)
            text_parts.append(f"[Page {i+1}]\n{processed_text}")
        return "\n\n".join(text_parts)
    
    def _process_pdf_text_for_tables(self, text: str) -> str:
        """
        Process PDF text to better preserve table-like structures.
        Identifies rows that look like table data and formats them for better indexing.
        """
        import re
        
        lines = text.split('\n')
        processed_lines = []
        
        id_pattern = re.compile(r'^([A-Z]{2,}-\d+)')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            id_match = id_pattern.match(line)
            if id_match:
                processed_lines.append(f"[ID: {id_match.group(1)}] {line}")
            else:
                processed_lines.append(line)
        
        return '\n'.join(processed_lines)
    
    def _parse_word_fallback(self, file_path: Path) -> str:
        """Parse Word document using python-docx."""
        from docx import Document as DocxDocument
        doc = DocxDocument(str(file_path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    
    def _parse_excel_fallback(self, file_path: Path) -> str:
        """Parse Excel file using openpyxl."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), data_only=True)
            content_parts = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                content_parts.append(f"[Sheet: {sheet_name}]")
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell else "" for cell in row])
                    if row_text.strip(" |"):
                        content_parts.append(row_text)
            return "\n".join(content_parts)
        except ImportError:
            return "Excel parsing requires openpyxl library"
    
    def _parse_ppt_fallback(self, file_path: Path) -> str:
        """Parse PowerPoint file using python-pptx."""
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            content_parts = []
            for i, slide in enumerate(prs.slides):
                content_parts.append(f"[Slide {i+1}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text)
            return "\n\n".join(content_parts)
        except ImportError:
            return "PowerPoint parsing requires python-pptx library"
    
    def _parse_txt_fallback(self, file_path: Path) -> str:
        """Parse text file."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    
    def _parse_image_fallback(self, file_path: Path) -> str:
        """Parse image file - returns placeholder for OCR."""
        return f"[Image file: {file_path.name}. OCR not available without Docling.]"
